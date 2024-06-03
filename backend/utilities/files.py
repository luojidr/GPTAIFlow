import json
import shutil
from pathlib import Path
import re
import pypdf
import mammoth
import openpyxl
from pptx import Presentation
from docx import Document
from utilities.print_utils import logger
import os
import subprocess
import tempfile
from typing import Optional, Union


def try_load_json_file(file_path: str, default=dict):
    if Path(file_path).exists():
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Failed to load json file: {file_path}, {e}")
    return default()


def get_files_contents(files: Union[str, list]):
    results = []
    if isinstance(files, str):
        files = [files]
    for file in files:
        try:
            if file.endswith(".docx"):
                doc = Document(file)
                text =''.join([paragraph.text for paragraph in doc.paragraphs])
                results.append(text)
                # with open(file, "rb") as docx_file:
                #     docx_data = mammoth.convert_to_markdown(docx_file)
                #     markdown_text = docx_data.value
                #     text_filter = re.findall(r'\(data.*?\)', markdown_text)
                #     if len(text_filter)>0 :
                #         logger.error(f"Failed to doc file: {file}")
                #         e =  ValueError('error file type ,include unknown file type ')
                #         e.error_detail = "doc文件包含非文字内容，请删除后重试"
                #         raise e
                #     results.append(markdown_text)
            elif file.endswith(".doc"):
                text = parse_doc(file)
                results.append(text)
            elif file.endswith(".pdf"):
                with open(file, "rb") as pdf_file_obj:
                    pdf_reader = pypdf.PdfReader(pdf_file_obj)
                    pdf_contents = [page.extract_text() for page in pdf_reader.pages]
                    results.append("\n\n".join(pdf_contents))
            elif file.endswith(".pptx"):
                ppt = Presentation(file)
                ppt_contents = []
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            ppt_contents.append(shape.text.strip())
                results.append("\n\n".join(ppt_contents))
            elif file.endswith(".xlsx"):
                wb = openpyxl.load_workbook(file, data_only=True)
                ws = wb.active
                csv_contents = []
                for row in ws.rows:
                    csv_contents.append(",".join([str(cell.value) if cell.value else "" for cell in row]))
                results.append("\n\n".join(csv_contents))
            elif file.endswith((".txt", ".md", ".html", ".json", ".csv")):
                try:
                    with open(file, "r", encoding="utf-8-sig") as txt_file:
                        txt_contents = txt_file.read()
                        results.append(txt_contents)
                except:
                    with open(file, "r", encoding="gbk") as txt_file:
                        txt_contents = txt_file.read()
                        results.append(txt_contents)
            elif file.endswith('.wav') or file.endswith(".mp3") or file.endswith(".m4a"):
                return files, "audio"
            else:
                raise ValueError("文件类型有误，请检查后重新上传")
                
        except Exception as e:
            if not hasattr(e,"e.error_detail"):
                e.error_detail = "文件无法正常打开，请确认选择无误且编码格式为utf-8"
            raise e

    return results, ""


def copy_file(src, dst):
    if Path(src).exists():
        Path(dst).parent.mkdir(parents=True, exist_ok=True)
        shutil.copy(src, dst)
        return True
    return False


def convert_office_doc(
    input_filename: str,
    output_directory: str,
    target_format: str = "docx",
    target_filter: Optional[str] = "MS Word 2007 XML",
):
    """Converts a .doc file to a .docx file using the libreoffice CLI.

    Parameters
    ----------
    input_filename: str
        The name of the .doc file to convert to .docx
    output_directory: str
        The output directory for the convert .docx file
    target_format: str
        The desired output format
    target_filter: str
        The output filter name to use when converting. See references below
        for details.
    """
    if target_filter is not None:
        target_format = f"{target_format}:{target_filter}"
    command = [
        "soffice",
        "--headless",
        "--convert-to",
        target_format,
        "--outdir",
        output_directory,
        input_filename,
    ]
    try:
        process = subprocess.Popen(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        output, error = process.communicate()
    except FileNotFoundError:
        raise FileNotFoundError(
            """soffice command was not found. Please install libreoffice
on your system and try again.
- Install instructions: https://www.libreoffice.org/get-help/install-howto/
- Mac: https://formulae.brew.sh/cask/libreoffice
- Debian: https://wiki.debian.org/LibreOffice""",
        )

    logger.info(output.decode().strip())
    if error:
        logger.error(error.decode().strip())


def parse_doc(file_path: str) -> str:
    with open(file_path, "rb") as file:
        tmp = tempfile.NamedTemporaryFile(delete=False)
        tmp.write(file.read())
        tmp.close()
        filename = tmp.name
        _, filename_no_path = os.path.split(os.path.abspath(tmp.name))
        base_filename, _ = os.path.splitext(filename_no_path)

        with tempfile.TemporaryDirectory() as tmpdir:
            convert_office_doc(
                filename,
                tmpdir,
                target_format="docx",
            )
            docx_filename = os.path.join(tmpdir, f"{base_filename}.docx")
            doc = Document(docx_filename)
            text = "".join([paragraph.text for paragraph in doc.paragraphs])
            if os.path.exists(docx_filename):
                os.unlink(docx_filename)

            return text